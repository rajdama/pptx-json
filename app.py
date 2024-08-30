from flask import Flask, request, send_file, jsonify
from pptx import Presentation
from io import BytesIO
from flask_cors import CORS 
from pptx.enum.lang import MSO_LANGUAGE_ID

app = Flask(__name__)
CORS(app) 

@app.route('/', methods=['POST'])
def update_ppt():
    if 'file' not in request.files:
        return jsonify({'error': 'No file or data provided'}), 400

    file = request.files['file']



    prs = Presentation(BytesIO(file.read()))
    
    slides = {}
    for slide_index, slide in enumerate(prs.slides):
        slides[f'slide{slide_index+1}'] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.language_id = MSO_LANGUAGE_ID.POLISH
                        txt = run.text.strip()
                        slides[f'slide{slide_index+1}'].append(txt)

    modified_ppt = BytesIO()
    prs.save(modified_ppt)
    modified_ppt.seek(0)

    return jsonify(slides)

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)
